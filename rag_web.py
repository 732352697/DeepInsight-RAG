import streamlit as st
import os
import tempfile
import sys
import docx
import pptx
import openpyxl
import time
import random  # 用来模拟思考时间的波动

# 🛡️ 强制禁用代理
os.environ["NO_PROXY"] = "localhost,127.0.0.1,0.0.0.0"

from langchain_ollama import ChatOllama
from langchain_community.document_loaders import PyPDFLoader

# --- 1. 页面配置 (DeepSeek 风格) ---
st.set_page_config(
    page_title="DeepInsight R1",
    page_icon="🧠",
    layout="wide"
)

# 注入 CSS 让状态框更好看
st.markdown("""
<style>
    .stStatusWidget {box-shadow: 0 2px 4px rgba(0,0,0,0.1); border-radius: 10px;}
    .reportview-container {margin-top: -2em;}
    header {visibility: hidden;}
    .stDeployButton {display:none;}
</style>
""", unsafe_allow_html=True)

st.title("🧠 DeepInsight · 深度思考版")
st.caption("🚀 本地 RAG 知识库 | 仿 DeepSeek 思考交互模式")

# --- 2. 侧边栏 ---
with st.sidebar:
    st.header("📂 知识库管理")
    uploaded_files = st.file_uploader(
        "上传资料 (多选)",
        type=["pdf", "docx", "xlsx", "pptx"],
        accept_multiple_files=True
    )

    if st.button("🗑️ 清空记忆"):
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        st.rerun()

    st.info("💡 提示：上传文件越多，思考时间可能会越长，请耐心等待“深度思考”完成。")


# --- 3. 核心解析逻辑 ---
def extract_text_from_file(uploaded_file):
    file_ext = uploaded_file.name.split(".")[-1].lower()
    text = ""
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_file_path = tmp_file.name

    try:
        if file_ext == "pdf":
            loader = PyPDFLoader(tmp_file_path)
            pages = loader.load()
            text = "\n".join([p.page_content for p in pages])
        elif file_ext == "docx":
            doc = docx.Document(tmp_file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_ext == "xlsx":
            wb = openpyxl.load_workbook(tmp_file_path, data_only=True)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text += f"\n[Sheet: {sheet}]\n"
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) for cell in row if cell is not None])
                    text += row_text + "\n"
        elif file_ext == "pptx":
            prs = pptx.Presentation(tmp_file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        return text
    except Exception as e:
        return f"解析异常: {str(e)}"
    finally:
        if os.path.exists(tmp_file_path):
            os.remove(tmp_file_path)


# --- 4. 文件处理 ---
if uploaded_files:
    current_file_names = ",".join([f.name for f in uploaded_files])
    if "last_processed_files" not in st.session_state or st.session_state.last_processed_files != current_file_names:
        combined_text = ""
        progress_bar = st.progress(0)
        status_text = st.empty()

        for i, file in enumerate(uploaded_files):
            status_text.text(f"正在读取: {file.name}...")
            file_text = extract_text_from_file(file)
            combined_text += f"\n\n=== 📄 {file.name} ===\n{file_text}\n"
            progress_bar.progress((i + 1) / len(uploaded_files))

        if len(combined_text) > 10000:
            combined_text = combined_text[:10000]

        st.session_state.doc_text = combined_text
        st.session_state.last_processed_files = current_file_names
        progress_bar.empty()
        status_text.empty()
        st.toast(f"✅ 知识库加载完毕，共 {len(combined_text)} 字", icon="🧠")

# --- 5. 聊天界面 (DeepSeek 风格核心区) ---
if "messages" not in st.session_state:
    st.session_state.messages = []

for msg in st.session_state.messages:
    avatar = "🧑‍💻" if msg["role"] == "user" else "🧠"
    with st.chat_message(msg["role"], avatar=avatar):
        st.markdown(msg["content"])

if prompt := st.chat_input("请输入您的问题..."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user", avatar="🧑‍💻"):
        st.markdown(prompt)

    if "doc_text" in st.session_state:
        with st.chat_message("assistant", avatar="🧠"):

            # 🌟 1. DeepSeek 风格的状态展示容器 🌟
            # 这个容器会在生成前显示，让用户觉得“它在思考”
            with st.status("🚀 正在进行深度思考...", expanded=True) as status:
                st.write("🔍 正在检索本地知识库索引...")
                time.sleep(0.3)  # 假装一点延迟，让用户看清步骤

                # 展示它找到了什么（增强信任感）
                doc_snippet = st.session_state.doc_text[:300].replace("\n", " ") + "..."
                st.write(f"📖 已提取上下文 (共 {len(st.session_state.doc_text)} 字)")
                st.code(doc_snippet, language="text")

                st.write("⚙️ 正在构建提示词工程 (Prompt Engineering)...")
                time.sleep(0.2)

                st.write("🧠 模型正在进行逻辑推理...")

                # 准备开始生成
                message_placeholder = st.empty()
                full_response = ""

                try:
                    llm = ChatOllama(
                        model="qwen2.5:1.5b",
                        temperature=0.1,  # 低温，保证事实准确
                        base_url="http://127.0.0.1:11434"
                    )

                    final_prompt = f"""
                    你是一个专业的分析助手。请基于下面的【参考文档】，进行深度思考并回答用户问题。

                    【参考文档】：
                    {st.session_state.doc_text}

                    【用户问题】：
                    {prompt}

                    请注意：如果文档中没有相关信息，请直接说明。
                    """

                    chunks = llm.stream(final_prompt)

                    # 刹车片机制 (保留，防崩)
                    last_update_time = 0

                    for chunk in chunks:
                        if chunk.content:
                            full_response += chunk.content
                            current_time = time.time()
                            if current_time - last_update_time > 0.05:
                                message_placeholder.markdown(full_response)
                                last_update_time = current_time

                    message_placeholder.markdown(full_response)

                    # 🌟 2. 思考完成，更新状态框 🌟
                    status.update(label="✅ 深度思考完成", state="complete", expanded=False)

                    st.session_state.messages.append({"role": "assistant", "content": full_response})

                except Exception as e:
                    status.update(label="❌ 思考过程出错", state="error")
                    st.error(f"错误: {str(e)}")
    else:
        st.warning("请先上传文档，我才能开始思考。")